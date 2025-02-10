import json
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Sample JSON data
data = {
    "company": "Example Corp",
    "Date": "2023-10-01",
    "Service": "VAPT",
    "companyNamePlaceholder": "{{companyName}}",
    "scopePlaceholder": "{{scope}}",
    "bulletPoints": [
        "Web-App PT (Blackbox)",
        "On-Prem Infrastructure PT (Gray Box) on 40 Critical Servers",
        "On-Prem Infrastructure PT (Gray Box) on 40 Critical Servers",
        "On-Prem Infrastructure PT (Gray Box) on 40 Critical Servers"
    ],
    "imagePath": "path/to/new/image.jpg"
}

# Load the PowerPoint template
template_path = "template_proposal.pptx"
output_path = "output_proposal.pptx"
presentation = Presentation(template_path)

# Function to replace placeholders in the presentation
def replace_placeholders(presentation, data):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        for key, value in data.items():
                            placeholder = f"{{{{{key}}}}}"
                            if placeholder in text:
                                text = text.replace(placeholder, value)
                        run.text = text

# Function to insert numbered bullet points into a specific text frame
def insert_bullet_points(slide, placeholder, bullet_points):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, "")
                        for i, bullet_point in enumerate(bullet_points, start=1):
                            p = text_frame.add_paragraph()
                            p.text = bullet_point
                            p.level = 0
                            p.bullet = False  # Ensure the paragraph is not a bullet point
                            p.number = i  # Set the numbering

# Function to replace an image in the first slide
def replace_image(slide, image_path):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(image_path, left, top, width, height)

# Replace placeholders with data
replace_placeholders(presentation, data)

# Insert numbered bullet points where the scope placeholder is found
for slide in presentation.slides:
    insert_bullet_points(slide, data["scopePlaceholder"], data["bulletPoints"])

# Replace image in the first slide
replace_image(presentation.slides[0], data["imagePath"])

# Check if the output file exists and remove it
if os.path.exists(output_path):
    os.remove(output_path)

# Save the modified presentation
presentation.save(output_path)

print(f"Presentation saved as {output_path}")