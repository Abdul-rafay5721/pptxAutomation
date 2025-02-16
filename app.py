from flask import Flask, render_template, request, redirect, url_for, send_file, after_this_request
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt  # Add this import
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.dml.color import RGBColor  # Add this import
from werkzeug.utils import secure_filename
from datetime import datetime
import tempfile

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'temp')
app.config['PROPOSALS_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'proposals')
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Print debugging information
print(f"Current working directory: {os.getcwd()}")
print(f"Base directory: {BASE_DIR}")

# Function to get available templates
def get_available_templates():
    templates = []
    proposals_dir = app.config['PROPOSALS_FOLDER']
    if os.path.exists(proposals_dir):
        for file in os.listdir(proposals_dir):
            if file.endswith('.pptx'):
                templates.append(file)
    return templates

# Function to replace placeholders in the presentation
def replace_placeholders(presentation, data):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        for key, value in data.items():
                            placeholder = f"{{{{{key}}}}}"
                            if placeholder in text:
                                text = text.replace(placeholder, value)
                        run.text = text

# Function to insert numbered bullet points into a specific text frame
def insert_bullet_points(slide, placeholder, bullet_points):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, "")
                        for i, bullet_point in enumerate(bullet_points, start=1):
                            p = text_frame.add_paragraph()
                            p.text = bullet_point
                            p.level = 0
                            p.bullet = False  # Ensure the paragraph is not a bullet point
                            p.number = i  # Set the numbering

# Function to insert deliverable bullets without amounts into a specific text frame or table cell
def insert_deliverable_bullets(slide, placeholder, deliverable_bullets):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if placeholder in cell.text:
                        cell.text = ""  # Clear the cell text completely
                        for i, deliverable in enumerate(deliverable_bullets, start=1):
                            p = cell.text_frame.add_paragraph()
                            p.text = f"{i}. {deliverable}"
                            p.level = 0  # Ensure level 0 for standard bullets
                            p.space_before = Pt(0)  # Remove space before paragraph
                            p.space_after = Pt(0)   # Remove space after paragraph
                        return  # Exit after replacing the placeholder
        elif shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder in run.text:
                        run.text = ""  # Clear the text completely
                        for i, deliverable in enumerate(deliverable_bullets, start=1):
                            p = text_frame.add_paragraph()
                            p.text = f"{i}. {deliverable}"
                            p.level = 0  # Standard bullet level
                            p.space_before = Pt(0)  # Remove space before paragraph
                            p.space_after = Pt(0)   # Remove space after paragraph
                        return  # Exit after replacing the placeholder

# Function to insert deliverables and amounts into specific placeholders within a table
def insert_deliverables_and_amounts(presentation, deliverables, amounts):
    placeholders = [f"{{{{deliverable{i}}}}}" for i in range(1, 8)]
    amount_placeholders = [f"{{{{amount{i}}}}}" for i in range(1, 8)]
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                deliverable_rows = []
                for i, row in enumerate(table.rows):
                    if i == 0:  # Skip header row
                        continue
                    if any(placeholders[j] in cell.text for j in range(len(placeholders)) for cell in row.cells):
                        deliverable_rows.append(row)
                for i, row in enumerate(deliverable_rows):
                    if i < len(deliverables):
                        for cell in row.cells:
                            if placeholders[i] in cell.text:
                                cell.text = cell.text.replace(placeholders[i], deliverables[i])
                            if amount_placeholders[i] in cell.text:
                                cell.text = cell.text.replace(amount_placeholders[i], amounts[i])
                    else:
                        table._tbl.remove(row._tr)  # Remove extra deliverable rows

# Function to replace an image placeholder in the presentation
def replace_image_placeholder(presentation, placeholder, image_path):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, "")
                            left = shape.left
                            top = shape.top
                            slide.shapes.add_picture(image_path, left, top, width=Inches(1.75), height=Inches(1.0))
                            return

# Dictionary to keep track of which template has which slide number, starting positions, width, and height for inserting the rectangle
TEMPLATE_SLIDE_MAP = {
    "VAPT proposal.pptx": {
        "slide_index": 38,  # Slide 39 (0-based index)
        "left": 1.15,  # Starting position from left in cm
        "top": 8.4,  # Starting position from top in cm
        "width_per_week": 10.4,  # Width per week in cm
        "height": 2  # Height in cm
    },
    "SAMA CSF Audit - Proposal.pptx": {
        "slide_index": 33,  # Slide 34 (0-based index)
        "left": 1.3,  # Starting position from left in cm
        "top": 7.8,  # Starting position from top in cm
        "width_per_week": 4,  # Width per week in cm
        "height": 2  # Height in cm
    },
    # Add more templates and their corresponding slide numbers, positions, width, and height here
}

# Function to add a timeline rectangle based on the number of weeks
def add_timeline_rectangle(presentation, slide_index, timeline_name, weeks, left_offset, top, width_per_week, height):
    if slide_index < len(presentation.slides):
        slide = presentation.slides[slide_index]
        left = Inches(left_offset / 2.54)  # Convert cm to inches
        top = Inches(top / 2.54)  # Convert cm to inches
        width = Inches(width_per_week * weeks / 2.54)  # Convert cm to inches
        height = Pt(height * 28.3465)  # Convert cm to points (1 cm = 28.3465 points)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(65, 105, 225)  # Fill color royal blue
        shape.line.color.rgb = RGBColor(255, 255, 255)  # Border color white
        text_frame = shape.text_frame
        text_frame.text = timeline_name
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center-align the text
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

@app.route('/')
def index():
    templates = get_available_templates()
    return render_template('index.html', templates=templates)

@app.route('/generate', methods=['POST'])
def generate():
    try:
        # Create temp directory with proper permissions if it doesn't exist
        if not os.path.exists(app.config['UPLOAD_FOLDER']):
            os.makedirs(app.config['UPLOAD_FOLDER'], mode=0o777)

        # Get selected template from proposals folder
        selected_template = request.form.get('template')
        template_path = os.path.join(app.config['PROPOSALS_FOLDER'], selected_template)
        
        if not os.path.exists(template_path):
            return f'Selected template not found: {selected_template}', 404

        data = {
            "company": request.form['company'],
            "Date": request.form['date'],
            "Service": request.form['service'],
            "activityPoints": request.form.getlist('activityPoints'),
            "companyNamePlaceholder": "{{companyName}}",
            "scopePlaceholder": "{{scope}}",
            "bulletPoints": request.form.getlist('bulletPoints'),
            "deliverableBullets": request.form.getlist('deliverableBullets'),
            "amounts": request.form.getlist('amounts'),
            "withoutVat": request.form['withoutVat'],
            "total": request.form['total'],
            "timelineNames": request.form.getlist('timelineNames'),
            "weeksList": [int(weeks) for weeks in request.form.getlist('weeks')]
        }

        presentation = Presentation(template_path)

        # Replace placeholders with data
        replace_placeholders(presentation, data)

        # Insert numbered bullet points where the scope placeholder is found
        for slide in presentation.slides:
            insert_bullet_points(slide, data["scopePlaceholder"], data["bulletPoints"])

        # Insert deliverable bullets where the deliverable placeholder is found
        for slide in presentation.slides:
            insert_deliverable_bullets(slide, "{{deliverable}}", data["deliverableBullets"])
            # Insert activities using the same function but with activity placeholder
            insert_deliverable_bullets(slide, "{{activity}}", data["activityPoints"])

        # Insert deliverables and amounts into specific placeholders within a table
        insert_deliverables_and_amounts(presentation, data["deliverableBullets"], data["amounts"])

        # Replace placeholders with data including totals
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if "{{withoutVat}}" in cell.text:
                                cell.text = cell.text.replace("{{withoutVat}}", data["withoutVat"])
                                # Set text color to white
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                            if "{{total}}" in cell.text:
                                cell.text = cell.text.replace("{{total}}", data["total"])
                                # Set text color to white
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.color.rgb = RGBColor(255, 255, 255)

        # Replace image placeholder with the uploaded image
        if 'image' in request.files:
            image_file = request.files['image']
            if image_file.filename != '':
                image_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(image_file.filename))
                image_file.save(image_path)
                replace_image_placeholder(presentation, "{{image}}", image_path)

        # Get template-specific settings
        template_settings = TEMPLATE_SLIDE_MAP[selected_template]
        slide_index = template_settings["slide_index"]
        left_offset = template_settings["left"]
        top = template_settings["top"]
        width_per_week = template_settings["width_per_week"]
        height = template_settings["height"]

        # Add timeline rectangles based on the number of weeks
        for timeline_name, weeks in zip(data["timelineNames"], data["weeksList"]):
            add_timeline_rectangle(presentation, slide_index, timeline_name, weeks, left_offset, top, width_per_week, height)
            left_offset += width_per_week * weeks  # Increment left offset by width per week

        # Generate a unique filename for the output
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        company_name = secure_filename(data['company'])
        output_filename = f"proposal_{company_name}_{timestamp}.pptx"
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)

        # Save the modified presentation
        presentation.save(output_path)
        os.chmod(output_path, 0o644)

        # Simplified cleanup without threading
        @after_this_request
        def cleanup(response):
            try:
                if os.path.exists(output_path):
                    os.remove(output_path)
                if 'image' in request.files and request.files['image'].filename != '':
                    image_path = os.path.join(app.config['UPLOAD_FOLDER'], 
                                            secure_filename(request.files['image'].filename))
                    if os.path.exists(image_path):
                        os.remove(image_path)
            except Exception as e:
                app.logger.error(f"Error in cleanup: {e}")
            return response

        # Check if file exists and is readable before sending
        if not os.path.exists(output_path):
            return "Failed to create output file", 500
            
        if not os.access(output_path, os.R_OK):
            return "Permission denied accessing output file", 500

        print(f"Sending file: {output_path}")  # Debug print
        return send_file(
            output_path,
            as_attachment=True,
            download_name=output_filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        app.logger.error(f"Error generating presentation: {e}")
        import traceback
        print(traceback.format_exc())  # Debug print
        return f"Error generating presentation: {str(e)}", 500

if __name__ == '__main__':
    # Development settings for Ubuntu
    app.run(host='0.0.0.0', port=8080, debug=True)
